#!/usr/bin/env python3
"""SOP FEAR 5 · La Querencia · Escalera V0–V3 (PPTX) con membrete Cacao Colab."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

FOREST = RGBColor(0x1A, 0x2E, 0x10)
CREAM = RGBColor(0xF7, 0xF1, 0xEE)
YELLOW = RGBColor(0xF2, 0xC8, 0x30)
CORAL = RGBColor(0xFF, 0x6A, 0x3D)
MUTED = RGBColor(0x5C, 0x6B, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
POD = RGBColor(0x86, 0xB6, 0x6B)

W = Inches(13.333)
H = Inches(7.5)
OUT = Path(__file__).resolve().parent / "SOP-V0-V3-FEAR5-La-Querencia.pptx"
OUT_LEGACY = Path(__file__).resolve().parent / "SOP-B0-V0-FEAR5-La-Querencia.pptx"
ART = Path("/opt/cursor/artifacts/SOP-V0-V3-FEAR5-La-Querencia.pptx")


def set_run(run, *, size=18, bold=False, color=FOREST, font="Georgia"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    try:
        rPr = run._r.get_or_add_rPr()
        latin = rPr.get_or_add_latin()
        latin.typeface = font
    except Exception:
        pass


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=FOREST, font="Georgia", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def letterhead(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST
    bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.72), W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = YELLOW
    accent.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(0.14), Inches(6), Inches(0.45), "CACAO COLAB", size=20, bold=True, color=YELLOW, font="Georgia")
    add_textbox(
        slide,
        Inches(6.8),
        Inches(0.18),
        Inches(5.9),
        Inches(0.4),
        "SOP · FEAR 5 · Escalera V0–V3  ·  Confidencial nodo",
        size=12,
        color=CREAM,
        font="Calibri",
        align=PP_ALIGN.RIGHT,
    )


def footer(slide, page: int, total: int):
    add_textbox(
        slide,
        Inches(0.55),
        Inches(7.05),
        Inches(9),
        Inches(0.3),
        "La Querencia · FEAR 5 · V0→V3  ·  cacaocolab.org",
        size=11,
        color=MUTED,
        font="Calibri",
    )
    add_textbox(slide, Inches(11.2), Inches(7.05), Inches(1.5), Inches(0.3), f"{page} / {total}", size=11, color=MUTED, font="Calibri", align=PP_ALIGN.RIGHT)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bullet_block(slide, left, top, width, height, lines, *, size=16, color=FOREST):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = "•  " + line
        set_run(run, size=size, color=color, font="Calibri")
    return box


def card(slide, left, top, width, height, fill=CREAM, line=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = RGBColor(0xD8, 0xCF, 0xC2)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def numbered_steps(slide, steps, *, start_top=1.7, left=0.55):
    for i, (num, title, body) in enumerate(steps):
        top = Inches(start_top + i * 0.88)
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), top, Inches(0.5), Inches(0.5))
        oval.fill.solid()
        oval.fill.fore_color.rgb = YELLOW
        oval.line.fill.background()
        add_textbox(slide, Inches(left), top + Inches(0.08), Inches(0.5), Inches(0.35), num, size=13, bold=True, color=FOREST, font="Calibri", align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(left + 0.7), top, Inches(11.2), Inches(0.32), title, size=17, bold=True, color=FOREST, font="Georgia")
        add_textbox(slide, Inches(left + 0.7), top + Inches(0.32), Inches(11.2), Inches(0.45), body, size=13, color=MUTED, font="Calibri")


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = 14

    # 1 Portada
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = FOREST
    bg.line.fill.background()
    yellow_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.55), W, Inches(0.95))
    yellow_bar.fill.solid()
    yellow_bar.fill.fore_color.rgb = YELLOW
    yellow_bar.line.fill.background()
    add_textbox(s, Inches(0.7), Inches(0.55), Inches(11), Inches(0.35), "CACAO COLAB  ×  LA QUERENCIA", size=14, bold=True, color=YELLOW, font="Calibri")
    add_textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(1.0), "SOP FEAR 5 · Escalera V0 → V3", size=40, bold=True, color=CREAM, font="Georgia")
    add_textbox(s, Inches(0.7), Inches(2.45), Inches(12), Inches(0.55), "Micro-lote → control térmico → tanque → bioreactor", size=22, color=YELLOW, font="Georgia")
    add_textbox(
        s,
        Inches(0.7),
        Inches(3.3),
        Inches(11.5),
        Inches(1.6),
        "Programa de experimentación en fermentación de precisión.\n"
        "Validamos curvas de temperatura (45 °C) y pH (≤ 4,5), y luego\n"
        "la acumulación / decaimiento de precursores de sabor del FEAR 5\n"
        "(referencia paper Santander 2025 · 7 biomarcadores).",
        size=17,
        color=CREAM,
        font="Calibri",
    )
    add_textbox(s, Inches(0.7), Inches(6.75), Inches(8), Inches(0.4), "Para: Rafael  ·  Operación de finca / fermentación", size=15, bold=True, color=FOREST, font="Calibri")
    add_textbox(s, Inches(9.0), Inches(6.75), Inches(3.7), Inches(0.4), "Documento SOP · Sep 2026", size=14, color=FOREST, font="Calibri", align=PP_ALIGN.RIGHT)

    # 2 Objetivo del programa
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.5), "1. Qué buscamos en todo el programa", size=28, bold=True, color=FOREST, font="Georgia")
    bullet_block(
        s,
        Inches(0.55),
        Inches(1.7),
        Inches(12),
        Inches(4.8),
        [
            "Material: FEAR 5 de La Querencia (mini-cosecha ~12 días; la fecha depende del clima).",
            "Hipótesis central: subir rápido a 45 °C, sostener, y bajar el pH lentamente hasta ≤ 4,5 en 72–96 h.",
            "Quitar ~20 % del mucílago con prensa de algodón; recoger en cantinas de miel; refrigerar/congelar.",
            "Medir T° y pH día a día; muestrear 500 g a las 72 h y 500 g a las 96 h.",
            "Agrupar los 7 biomarcadores precursores de sabor (paper Santander 2025) y ver su decaimiento",
            "    cuando hay error de fermentación o pérdida al pasar de V0→V1→V2→V3.",
            "Escala de masa: 1,5 kg (V0/V1) → ~15 kg (V2) → ~150 kg (V3 bioreactor tipo Agrosavia).",
        ],
        size=17,
    )
    footer(s, 2, total)

    # 3 Mapa V0-V3
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "2. Mapa de versiones", size=28, bold=True, color=FOREST, font="Georgia")
    add_textbox(s, Inches(0.55), Inches(1.55), Inches(12), Inches(0.35), "Cada versión valida una cosa. No saltar etapas.", size=15, color=MUTED, font="Calibri")

    stages = [
        ("V0", "B0 · 1,5 kg", "Tanque + cesta\nTermófilas / sol\nControl manual\nT° vs pH", YELLOW),
        ("V1", "B1 · 1,5–15 kg", "Bolsas al vacío\nen el mismo tanque\n+ retrocirculador\nde agua a 45 °C", CREAM),
        ("V2", "B2 · ~15 kg", "Tanque brewing /\nyogur / amasado\nTemp. constante\n45 °C", CREAM),
        ("V3", "B3 · ~150 kg", "Bioreactor\nestilo Agrosavia\nEscalamiento\ntecnificado", CREAM),
    ]
    for i, (ver, sub, body, fill) in enumerate(stages):
        left = Inches(0.55 + i * 3.15)
        card(s, left, Inches(2.1), Inches(2.95), Inches(4.2), fill=fill)
        add_textbox(s, left + Inches(0.15), Inches(2.3), Inches(2.65), Inches(0.45), ver, size=28, bold=True, color=FOREST, font="Georgia")
        add_textbox(s, left + Inches(0.15), Inches(2.85), Inches(2.65), Inches(0.4), sub, size=13, bold=True, color=CORAL if fill == YELLOW else MUTED, font="Calibri")
        add_textbox(s, left + Inches(0.15), Inches(3.4), Inches(2.65), Inches(2.5), body, size=14, color=FOREST, font="Calibri")
    footer(s, 3, total)

    # 4 V0 qué valida
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.5), "3. V0 · Qué validamos (solo esto)", size=28, bold=True, color=FOREST, font="Georgia")
    card(s, Inches(0.55), Inches(1.75), Inches(12.2), Inches(1.3), YELLOW)
    add_textbox(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(11.7),
        Inches(0.9),
        "Assumption V0: la temperatura puede subir rápido a 45 °C y sostenerse,\n"
        "mientras el pH baja lento hasta ≤ 4,5 entre 72 y 96 horas — sin bioreactor.",
        size=17,
        bold=True,
        color=FOREST,
        font="Calibri",
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(3.3),
        Inches(12),
        Inches(3.2),
        [
            "Montaje: tanque de acero + cesta con 1,5 kg de grano FEAR 5, sin agua, tapado.",
            "Calor: bacterias termófilas. Si no alcanza 45 °C → ayudar al sol (aeróbico).",
            "0–48 h aeróbico (arranque). Luego anaeróbico tapado hasta 72–96 h.",
            "Opción: pasar grano desde el cajón cuando ya tenga calor hacia el tanque.",
            "Éxito V0 = bitácora T°/pH + muestras 72 h y 96 h. Sin V0 no hay V1.",
        ],
        size=17,
    )
    footer(s, 4, total)

    # 5 V0 materiales + prep
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "4. V0 · Materiales y preparación", size=28, bold=True, color=FOREST, font="Georgia")
    card(s, Inches(0.55), Inches(1.65), Inches(6.0), Inches(4.7), CREAM)
    add_textbox(s, Inches(0.75), Inches(1.85), Inches(5.5), Inches(0.35), "Equipo", size=17, bold=True, color=FOREST, font="Georgia")
    bullet_block(
        s,
        Inches(0.75),
        Inches(2.3),
        Inches(5.5),
        Inches(3.8),
        [
            "Tanque acero inoxidable limpio",
            "Cesta (grano sin agua)",
            "Tapa / cobertor anaeróbico",
            "Termómetro + pH-metro calibrado",
            "Prensa de algodón manual",
            "Cantinas de miel (mucílago)",
            "Balanza + bolsas 500 g",
            "Bitácora / libreta",
        ],
        size=15,
    )
    card(s, Inches(6.8), Inches(1.65), Inches(5.9), Inches(4.7), CREAM)
    add_textbox(s, Inches(7.0), Inches(1.85), Inches(5.5), Inches(0.35), "Día 0 · Preparación", size=17, bold=True, color=FOREST, font="Georgia")
    bullet_block(
        s,
        Inches(7.0),
        Inches(2.3),
        Inches(5.5),
        Inches(3.8),
        [
            "Cosechar FEAR 5 maduro",
            "Quitar ~20 % mucílago (prensa)",
            "No lavar el grano con agua",
            "Mucílago → cantinas → frío",
            "Cargar 1,5 kg en cesta",
            "Cubrir cotiledón / masa",
            "Arranque térmico (termófilas/sol)",
            "Anotar lote, fecha, kilos",
        ],
        size=15,
    )
    footer(s, 5, total)

    # 6 V0 protocolo fermentación
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "5. V0 · Protocolo de fermentación", size=28, bold=True, color=FOREST, font="Georgia")
    card(s, Inches(0.55), Inches(1.65), Inches(6.0), Inches(4.7), CREAM)
    add_textbox(s, Inches(0.75), Inches(1.85), Inches(5.5), Inches(0.4), "Fase A · Aeróbica 0–48 h", size=18, bold=True, color=FOREST, font="Georgia")
    bullet_block(
        s,
        Inches(0.75),
        Inches(2.4),
        Inches(5.5),
        Inches(3.6),
        [
            "Grano en cesta, sin agua",
            "Meta: subir rápido a 45 °C",
            "Si no sube: sol controlado",
            "T° mínimo 2×/día",
            "pH 1×/día (misma hora)",
            "No agregar agua",
        ],
        size=16,
    )
    card(s, Inches(6.8), Inches(1.65), Inches(5.9), Inches(4.7), FOREST)
    add_textbox(s, Inches(7.0), Inches(1.85), Inches(5.5), Inches(0.4), "Fase B · Anaeróbica 48–96 h", size=18, bold=True, color=YELLOW, font="Georgia")
    bullet_block(
        s,
        Inches(7.0),
        Inches(2.4),
        Inches(5.5),
        Inches(3.6),
        [
            "Tapar bien (anaeróbico)",
            "Sostener ~45 °C constantes",
            "pH baja lento → ≤ 4,5",
            "Seguir bitácora diaria",
            "72 h → Muestra A (500 g)",
            "96 h → Muestra B (500 g)",
        ],
        size=16,
        color=CREAM,
    )
    footer(s, 6, total)

    # 7 V0 checklist
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "6. V0 · Checklist de cierre", size=28, bold=True, color=FOREST, font="Georgia")
    numbered_steps(
        s,
        [
            ("01", "Bitácora completa", "T° (≥2/día) y pH (1/día) de 0 a 96 h, con iniciales."),
            ("02", "Curva de temperatura", "Llegó a ~45 °C y se sostuvo la mayor parte de las 72 h."),
            ("03", "Curva de pH", "Bajó gradual hasta ≤ 4,5 — o se documentó por qué no."),
            ("04", "Dos muestras", "500 g a 72 h y 500 g a 96 h, etiquetadas y conservadas."),
            ("05", "Mucílago guardado", "Cantinas etiquetadas (fecha · FEAR 5 · Querencia) en frío."),
            ("06", "Fotos", "Montaje inicial (cesta/tanque) y cierre. Listo para decidir V1."),
        ],
        start_top=1.6,
    )
    footer(s, 7, total)

    # 8 V1 qué valida
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "7. V1 · Temperatura constante con agua", size=28, bold=True, color=FOREST, font="Georgia")
    card(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(1.2), YELLOW)
    add_textbox(
        s,
        Inches(0.8),
        Inches(1.95),
        Inches(11.7),
        Inches(0.8),
        "Assumption V1: con grano en bolsas al vacío dentro del mismo tanque,\n"
        "un retrocirculador de agua a 45 °C sostiene la temperatura sin depender del sol.",
        size=16,
        bold=True,
        color=FOREST,
        font="Calibri",
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(3.15),
        Inches(12),
        Inches(3.4),
        [
            "Solo después de V0 aprobado (curvas entendidas).",
            "Grano FEAR 5 en bolsas al vacío (sub-B) dentro del tanque de acero.",
            "Agua recirculando alrededor de las bolsas a setpoint 45 °C.",
            "Misma lógica de pH (descenso lento a ≤ 4,5) y muestreo 72 h / 96 h.",
            "Comparamos vs V0: ¿menos variación de T°? ¿mejor o peor precursor?",
            "Escala aún micro / piloto (1,5–15 kg según capacidad del tanque).",
        ],
        size=16,
    )
    footer(s, 8, total)

    # 9 V1 paso a paso
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "8. V1 · Paso a paso para Rafael", size=28, bold=True, color=FOREST, font="Georgia")
    numbered_steps(
        s,
        [
            ("01", "Preparar igual que V0", "Mucílago ~20 %, cantinas, etiquetas. Grano listo y limpio."),
            ("02", "Envasar al vacío", "Porciones en bolsas al vacío (sub-B). Sin agua libre dentro de la bolsa."),
            ("03", "Montar baño", "Bolsas dentro del tanque. Agua alrededor. Retrocirculador a 45 °C."),
            ("04", "Correr 72–96 h", "Registrar T° del agua/grano y pH del lote según protocolo acordado."),
            ("05", "Muestrear", "500 g a 72 h y 500 g a 96 h. Comparar con las muestras V0."),
            ("06", "Decidir V2", "Si T° es estable y pH llega: pasar a tanque brewing/yogur (~15 kg)."),
        ],
        start_top=1.55,
    )
    footer(s, 9, total)

    # 10 V2
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "9. V2 · Tanque tipo brewing / yogur", size=28, bold=True, color=FOREST, font="Georgia")
    card(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(1.2), YELLOW)
    add_textbox(
        s,
        Inches(0.8),
        Inches(1.95),
        Inches(11.7),
        Inches(0.8),
        "Assumption V2: un tanque de cervecería, amasado o fermentación de yogur\n"
        "mantiene 45 °C constante a escala ~15 kg, acercándonos al bioreactor.",
        size=16,
        bold=True,
        color=FOREST,
        font="Calibri",
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(3.15),
        Inches(12),
        Inches(3.4),
        [
            "Equipo: tanque con control térmico (jacket / resistencia / baño) a 45 °C.",
            "Carga objetivo ~15 kg de FEAR 5 preparado (misma lógica de mucílago 20 %).",
            "Protocolo de fases: arranque aeróbico corto si aplica → anaeróbico tapado.",
            "Misma bitácora T°/pH; mismas ventanas 72 h y 96 h para muestras.",
            "Pregunta clave: ¿se conservan los precursores al subir de 1,5 kg a 15 kg?",
            "Documentar cualquier error de fermentación o pérdida por transferencia.",
        ],
        size=16,
    )
    footer(s, 10, total)

    # 11 V3
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "10. V3 · Bioreactor ~150 kg (Agrosavia-style)", size=26, bold=True, color=FOREST, font="Georgia")
    card(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(1.2), YELLOW)
    add_textbox(
        s,
        Inches(0.8),
        Inches(1.95),
        Inches(11.7),
        Inches(0.8),
        "Assumption V3: el protocolo ya validado en V0–V2 escala a un bioreactor\n"
        "tecnificado (~150 kg) sin perder el perfil de precursores del FEAR 5.",
        size=16,
        bold=True,
        color=FOREST,
        font="Calibri",
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(3.15),
        Inches(12),
        Inches(3.4),
        [
            "Equipo: bioreactor / fermentador tecnificado (referencia tipo Agrosavia).",
            "Setpoint 45 °C controlado; monitoreo continuo de T° y pH.",
            "SOP de planta: limpieza, carga, muestreo, trazabilidad de lote.",
            "Comparar biomarcadores V0→V1→V2→V3: acumulación vs decaimiento.",
            "Si hay pérdida: ¿error de fermentación o transferencia tecnológica?",
            "V3 es escalamiento, no inventar de cero: hereda setpoints de V0–V2.",
        ],
        size=16,
    )
    footer(s, 11, total)

    # 12 Biomarcadores
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "11. Precursores · qué medimos entre versiones", size=26, bold=True, color=FOREST, font="Georgia")
    bullet_block(
        s,
        Inches(0.55),
        Inches(1.7),
        Inches(12),
        Inches(2.4),
        [
            "Referencia: paper Santander 2025 · FEAR 5 · 7 biomarcadores precursores de sabor.",
            "En V0 solo acumulamos evidencia de curvas + muestras 72/96 h.",
            "Desde V1 comparamos: misma genética, distinto control térmico.",
            "Buscamos agrupar precursores y detectar decaimiento por etapa (V0→V3).",
        ],
        size=16,
    )
    # comparison table header cards
    headers = [("Versión", "Escala", "Control T°", "Pregunta de validación")]
    rows = [
        ("V0", "1,5 kg", "Termófilas / sol", "¿Sube a 45 °C y baja pH?"),
        ("V1", "1,5–15 kg", "Agua 45 °C + vacío", "¿T° más estable que V0?"),
        ("V2", "~15 kg", "Tanque brewing/yogur", "¿Precursores aguantan escala?"),
        ("V3", "~150 kg", "Bioreactor", "¿Se escala sin decaer?"),
    ]
    table = s.shapes.add_table(5, 4, Inches(0.55), Inches(4.15), Inches(12.2), Inches(2.4)).table
    widths = [1.8, 2.2, 3.6, 4.6]
    for j, w in enumerate(widths):
        table.columns[j].width = Inches(w)
    for j, htxt in enumerate(["Versión", "Escala", "Control T°", "Pregunta de validación"]):
        cell = table.cell(0, j)
        cell.text = htxt
        cell.fill.solid()
        cell.fill.fore_color.rgb = YELLOW
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                set_run(run, size=13, bold=True, color=FOREST, font="Calibri")
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run(run, size=12, bold=(j == 0), color=FOREST, font="Calibri")
    footer(s, 12, total)

    # 13 Bitácora
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "12. Bitácora común (todas las versiones)", size=26, bold=True, color=FOREST, font="Georgia")
    rows_n, cols_n = 6, 7
    table = s.shapes.add_table(rows_n, cols_n, Inches(0.4), Inches(1.7), Inches(12.5), Inches(4.5)).table
    headers = ["Versión", "Hora", "T °C", "pH", "Fase", "Observación", "Iniciales"]
    widths = [1.4, 1.4, 1.3, 1.3, 1.5, 4.0, 1.6]
    for j, htxt in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = htxt
        cell.fill.solid()
        cell.fill.fore_color.rgb = YELLOW
        table.columns[j].width = Inches(widths[j])
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                set_run(run, size=12, bold=True, color=FOREST, font="Calibri")
    examples = [
        ["V0", "0 h", "", "", "A", "Carga 1,5 kg / −20% mucílago", ""],
        ["V0", "48 h", "", "", "A→B", "Cierre anaeróbico", ""],
        ["V0", "72 h", "", "", "B", "Muestra A 500 g", ""],
        ["V0", "96 h", "", "", "B", "Muestra B 500 g", ""],
        ["V1…", "", "", "", "", "Repetir filas por versión", ""],
    ]
    for i, row in enumerate(examples, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run(run, size=11, color=FOREST, font="Calibri")
    footer(s, 13, total)

    # 14 Cierre
    s = blank_slide(prs)
    letterhead(s)
    add_textbox(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.45), "13. Orden de trabajo para Rafael", size=28, bold=True, color=FOREST, font="Georgia")
    numbered_steps(
        s,
        [
            ("V0", "Ahora", "1,5 kg · tanque + cesta · termófilas/sol · curvas T°/pH · muestras 72/96 h."),
            ("V1", "Después de V0", "Bolsas al vacío + retrocirculador de agua a 45 °C en el mismo tanque."),
            ("V2", "Después de V1", "Tanque brewing/yogur ~15 kg con temperatura constante."),
            ("V3", "Después de V2", "Bioreactor ~150 kg · escalar sin perder precursores FEAR 5."),
        ],
        start_top=1.7,
    )
    card(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.2), FOREST)
    add_textbox(
        s,
        Inches(0.8),
        Inches(5.65),
        Inches(11.7),
        Inches(0.8),
        "Regla Colab: una versión a la vez. Rafael ejecuta. El Colab interpreta curvas y precursores.\n"
        "Mini-cosecha FEAR 5 en ~12 días (clima manda). Arrancamos en V0.",
        size=15,
        color=CREAM,
        font="Calibri",
    )
    footer(s, 14, total)

    prs.save(OUT)
    prs.save(OUT_LEGACY)  # mantener nombre previo accesible
    prs.save(ART)
    print(f"Wrote {OUT}")
    print(f"Wrote {OUT_LEGACY}")
    print(f"Wrote {ART}")


if __name__ == "__main__":
    build()
